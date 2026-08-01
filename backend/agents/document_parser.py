import os
import re
import fitz
from docx import Document
from pptx import Presentation
from backend.schemas.models import DocumentStructure, Section, TableData


def parse_pdf(file_path):
    sections = []
    tables = []
    equations = []
    full_text = ""

    with fitz.open(file_path) as doc:
        for page_num in range(len(doc)):
            page = doc[page_num]
            blocks = page.get_text("dict")["blocks"]

            for block in blocks:
                if block["type"] == 0:
                    block_text = ""
                    max_font_size = 0

                    for line in block["lines"]:
                        for span in line["spans"]:
                            block_text += span["text"] + " "
                            if span["size"] > max_font_size:
                                max_font_size = span["size"]

                    block_text = block_text.strip()
                    if not block_text:
                        continue

                    level = 0
                    if max_font_size > 16:
                        level = 1
                    elif max_font_size > 13:
                        level = 2
                    elif max_font_size > 11:
                        level = 3

                    if level > 0:
                        sections.append(Section(
                            heading=block_text,
                            content="",
                            level=level,
                            page_number=page_num + 1
                        ))
                    else:
                        if sections:
                            sections[-1].content += block_text + "\n"

                    full_text += block_text + "\n"

                    math_patterns = re.findall(r'\$[^$]+\$|[A-Za-z]\s*=\s*[^,\n]+', block_text)
                    equations.extend(math_patterns)

            page_tables = page.find_tables()
            for table in page_tables:
                table_data = table.extract()
                if table_data and len(table_data) > 1:
                    headers = [str(h) for h in table_data[0]] if table_data[0] else []
                    rows = [[str(cell) for cell in row] for row in table_data[1:]]
                    tables.append(TableData(
                        title=f"Table on page {page_num + 1}",
                        headers=headers,
                        rows=rows
                    ))

    if not sections:
        paragraphs = full_text.split("\n\n")
        for i, para in enumerate(paragraphs):
            if para.strip():
                sections.append(Section(
                    heading=f"Section {i+1}",
                    content=para.strip(),
                    level=0,
                    page_number=1
                ))

    return full_text, sections, tables, equations


def parse_docx(file_path):
    doc = Document(file_path)
    sections = []
    tables = []
    full_text = ""
    current_section = None

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        full_text += text + "\n"

        if para.style.name.startswith("Heading"):
            try:
                level = int(para.style.name.replace("Heading ", ""))
            except ValueError:
                level = 1

            current_section = Section(
                heading=text,
                content="",
                level=level,
                page_number=0
            )
            sections.append(current_section)
        else:
            if current_section:
                current_section.content += text + "\n"
            else:
                current_section = Section(
                    heading="Introduction",
                    content=text + "\n",
                    level=0,
                    page_number=0
                )
                sections.append(current_section)

    for table in doc.tables:
        headers = [cell.text for cell in table.rows[0].cells] if table.rows else []
        rows = []
        for row in table.rows[1:]:
            rows.append([cell.text for cell in row.cells])
        tables.append(TableData(title="", headers=headers, rows=rows))

    return full_text, sections, tables, []


def parse_pptx(file_path):
    prs = Presentation(file_path)
    sections = []
    full_text = ""

    for slide_num, slide in enumerate(prs.slides):
        slide_text = ""
        title = f"Slide {slide_num + 1}"

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_text += text + "\n"

            if hasattr(shape, "title") and shape == slide.shapes.title:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip()

        if slide_text.strip():
            sections.append(Section(
                heading=title,
                content=slide_text,
                level=1,
                page_number=slide_num + 1
            ))
            full_text += slide_text + "\n"

    return full_text, sections, [], []


def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8") as f:
        text = f.read()

    sections = []
    paragraphs = text.split("\n\n")

    for i, para in enumerate(paragraphs):
        para = para.strip()
        if not para:
            continue

        lines = para.split("\n")
        heading = lines[0] if len(lines[0]) < 100 else f"Section {i+1}"
        content = "\n".join(lines[1:]) if len(lines) > 1 else para

        sections.append(Section(
            heading=heading,
            content=content,
            level=0,
            page_number=0
        ))

    return text, sections, [], []


def parse_document(file_path):
    filename = os.path.basename(file_path)
    extension = os.path.splitext(filename)[1].lower()

    if extension == ".pdf":
        full_text, sections, tables, equations = parse_pdf(file_path)
        file_type = "PDF"
    elif extension == ".docx":
        full_text, sections, tables, equations = parse_docx(file_path)
        file_type = "DOCX"
    elif extension in [".pptx", ".ppt"]:
        full_text, sections, tables, equations = parse_pptx(file_path)
        file_type = "PPT"
    elif extension == ".txt":
        full_text, sections, tables, equations = parse_txt(file_path)
        file_type = "TXT"
    else:
        raise ValueError(f"Unsupported file type: {extension}")

    word_count = len(full_text.split())

    return DocumentStructure(
        filename=filename,
        file_type=file_type,
        total_pages=len(sections),
        word_count=word_count,
        raw_text=full_text[:300000],
        sections=sections,
        tables=tables,
        equations=equations,
        metadata={"source_path": file_path}
    )
